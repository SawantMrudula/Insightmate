import streamlit as st
import pandas as pd
import pdfplumber
import docx
from pptx import Presentation
from io import BytesIO
import matplotlib.pyplot as plt
import seaborn as sns
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain.schema import HumanMessage
from dotenv import load_dotenv
import os
from io import BytesIO
from PyPDF2 import PdfReader
from tenacity import retry, stop_after_attempt, wait_fixed

# Load environment variables from .env file
load_dotenv()

# Set environment variables
os.environ['GOOGLE_API_KEY'] = os.getenv("GOOGLE_API_KEY")
os.environ["LANGSMITH_TRACING"] = "true"
os.environ["LANGCHAIN_API_KEY"] = os.getenv("LANGCHAIN_API_KEY")
MAX_TOKENS = 50

# Initialize session state for page navigation
if "page" not in st.session_state:
    st.session_state.page = 'main'

# Set a theme for plots using Seaborn
sns.set_theme(style="whitegrid", palette="pastel")

# Initialize the LangChain model
chat_llm = ChatGoogleGenerativeAI(model="gemini-1.5-pro")

# Prompt template for the model
prompt_template = ChatPromptTemplate.from_template(
    "You are an assistant knowledgeable in data analysis. Answer the following question based on the dataset: {question}")


def load_csv(file):
    return pd.read_csv(file)

def load_pdf(file):
    data = []
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                df = pd.DataFrame(table[1:], columns=table[0])
                data.append(df)
    return pd.concat(data, ignore_index=True) if data else pd.DataFrame()

def load_docx(file):
    doc = docx.Document(file)
    data = []
    for table in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        df = pd.DataFrame(rows[1:], columns=rows[0])
        data.append(df)
    return pd.concat(data, ignore_index=True) if data else pd.DataFrame()

def load_pptx(file):
    presentation = Presentation(file)
    data = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows = [[cell.text_frame.text for cell in row.cells] for row in table.rows]
                df = pd.DataFrame(rows[1:], columns=rows[0])
                data.append(df)
    return pd.concat(data, ignore_index=True) if data else pd.DataFrame()

# Plot Functions
def plot_data(df, column):
    fig, ax = plt.subplots()
    sns.histplot(df[column], kde=True, ax=ax)
    return fig

# Generate AI Description

# def generate_ai_description(prompt):
#     chat_llm = ChatOpenAI()
#     messages = [HumanMessage(content=prompt)]
#     response = chat_llm(messages=messages)
#     return response.content

# File Upload and Processing
st.title("InsightMate: Data Analysis for Various File Formats")
file = st.file_uploader("Upload your file (CSV, PDF, DOCX, PPTX)", type=["csv", "pdf", "docx", "pptx"])

if file:
    file_extension = os.path.splitext(file.name)[1].lower()
    if file_extension == ".csv":
        df = load_csv(file)
    elif file_extension == ".pdf":
        df = load_pdf(file)
    elif file_extension == ".docx":
        df = load_docx(file)
    elif file_extension == ".pptx":
        df = load_pptx(file)
    else:
        st.error("Unsupported file format")
        st.stop()
    
    st.write("Data Preview:")
    st.write(df.head())
    
    # Column Selection for Visualization
    numeric_columns = df.select_dtypes(include=['float', 'int']).columns.tolist()
    categorical_columns = df.select_dtypes(include=['object']).columns.tolist()

    # Visualization for Numeric Columns
    if numeric_columns:
        st.write("### Numeric Column Visualizations")
        for column in numeric_columns:
            st.write(f"**{column}**")
            fig = plot_data(df, column)
            st.pyplot(fig)

            # AI Description
            description_prompt = f"Provide a simple explanation for the distribution of '{column}'. Max 45 words."
            ai_description = generate_ai_description(description_prompt)
            st.write(f"*AI's Description:* {ai_description}")

    # Visualization for Categorical Columns
    if categorical_columns:
        st.write("### Categorical Column Visualizations")
        for column in categorical_columns:
            fig, ax = plt.subplots()
            sns.countplot(data=df, x=column, ax=ax)
            plt.xticks(rotation=45)
            st.pyplot(fig)

            # AI Description
            description_prompt = f"Summarize the distribution of '{column}'. Max 45 words."
            ai_description = generate_ai_description(description_prompt)
            st.write(f"*AI's Description:* {ai_description}")

    # AI-Based Correlation Insights
    if len(numeric_columns) > 1:
        st.write("## AI-Based Correlation Insights")
        correlation_matrix = df[numeric_columns].corr().to_dict()
        if st.button("Generate AI Insights on Correlations"):
            formatted_corr_matrix = pd.DataFrame(correlation_matrix).to_string()
            ai_response = generate_ai_description(formatted_corr_matrix)
            st.write(f"AI Insights on Correlations: {ai_response}")
    else:
        st.write("Not enough numerical columns to compute correlations.")

# AI-Based Querying of Dataset
st.write("## Ask a Question about the Dataset")
user_question = st.text_input("Ask a question based on the data")

if user_question and 'df' in locals():
    data_context = df.head().to_string()
    prompt = f"Using the following data:\n\n{data_context}\n\nQuestion: {user_question}"
    ai_response = generate_ai_description(prompt)
    st.write(f"AI's Answer: {ai_response}")
    
# CSS to enhance UI (background & text style)
st.markdown("""
    <style>
        /* Main background gradient */
        .stApp {
            background-image: linear-gradient(to right, #28313b, #485461);
            font-family: 'Arial', sans-serif;
        }

        /* Headings color */
        h1, h2, h3 {
            color: #f1f1f1;  /* Light gray for a good contrast with the dark background */
        }

        /* Sidebar styles */
        .sidebar-content {
            background-color: #1f2a30;
            color: #ffffff;  /* Make the sidebar text white for visibility */
        }

        .sidebar .sidebar-content {
            padding: 10px;
        }

        /* Footer hidden */
        footer {
            visibility: hidden;
        }

        /* Container padding */
        .block-container {
            padding: 2rem;
        }

        /* Button style: Darker background with a lighter border and text */
        .stButton>button {
            background-color: #4e5a63;  /* Button background to match the gradient theme */
            color: #ffffff;  /* White text for good contrast */
            border: 2px solid #b0c4de;  /* Lighter blue-gray border for a subtle highlight */
            border-radius: 5px;
            padding: 10px 20px;
        }

        /* Input field style: Background matching the theme, text lighter */
        .stTextInput>div>div>input {
            background-color: #3b444b;
            color: #ffffff;  /* White text inside input fields */
            border: 2px solid #b0c4de;
        }

        /* Optional: Style for file uploader to match the theme */
        .stFileUploader>label {
            color: #ffffff;  /* White text for file upload label */
        }
    </style>
""", unsafe_allow_html=True)

# Title of the Streamlit app
st.title("DataViz: Explore Your Data with AI Insights")

# Sidebar for additional options and links
st.sidebar.title("Explore Options")
st.sidebar.markdown("Use the sidebar for additional actions or to navigate through different sections.")

# Initialize an empty list to store Q&A history
if 'qna_history' not in st.session_state:
    st.session_state.qna_history = []

# Navigation using st.session_state to toggle between main app and Q&A history page
if 'page' not in st.session_state:
    st.session_state.page = 'main'


# Function to display Q&A history page
def show_qna_history():
    st.write("## Q&A History")
    if len(st.session_state.qna_history) > 0:
        for i, entry in enumerate(st.session_state.qna_history, 1):
            st.write(f"### Question {i}:")
            st.write(f"Q: {entry['question']}")
            st.write(f"A: {entry['answer']}")
    else:
        st.write("No Q&A history available yet.")
    if st.button("Back to Main"):
        st.session_state.page = 'main'
        st.session_state.page = 'qna_history'


# Function to download a plot as a PNG file
def download_plot(fig, filename):
    buf = BytesIO()
    fig.savefig(buf, format="png")
    buf.seek(0)
    return buf


# Function to display main app page
def show_main_page():
    # Button to go to Q&A history page
    if st.sidebar.button("View Q&A History"):
        st.session_state.page = 'qna_history'
        st.experimental_rerun()

    # File uploader for user to upload CSV files
    uploaded_file = st.file_uploader("Upload a CSV, PDF, DOCX, or PPTX file", type=["csv", "pdf", "docx", "pptx"])

    if uploaded_file:
        if uploaded_file.name.endswith(".csv"):
            df = parse_file(uploaded_file)
            if isinstance(df, pd.DataFrame):
                st.write("### Dataset Preview")
                st.dataframe(df.head())


        # Display a preview of the dataset
        st.write("## Dataset Preview")
        st.dataframe(df.head())

        # Display basic dataset information
        st.write("## Basic Information about the Dataset")
        st.write(f"Number of rows: {df.shape[0]}")
        st.write(f"Number of columns: {df.shape[1]}")
        st.write("### Dataset Summary")
        st.write(df.describe())

        # Check for missing values
        if df.isnull().values.any():
            st.warning("This dataset contains missing values. Cleaning the dataset...")
            changes_made = []

            # Fill missing values in numeric columns with mean
            for col in df.select_dtypes(include=['float64', 'int64']).columns:
                if df[col].isnull().sum() > 0:
                    df[col].fillna(df[col].mean(), inplace=True)
                    changes_made.append(f"Filled missing values in numeric column '{col}' with mean value.")

            # Fill missing values in categorical columns with mode
            for col in df.select_dtypes(include=['object']).columns:
                if df[col].isnull().sum() > 0:
                    df[col].fillna(df[col].mode()[0], inplace=True)
                    changes_made.append(f"Filled missing values in categorical column '{col}' with mode.")

            # Display the preprocessing summary
            st.write("## Data Preprocessing Summary")
            if changes_made:
                for change in changes_made:
                    st.write(f"- {change}")
            else:
                st.write("No additional cleaning was necessary.")
        else:
            st.success("No missing or inconsistent values found in the dataset.")

        # Identify numerical and categorical columns
        numeric_columns = df.select_dtypes(include=['float64', 'int64']).columns
        categorical_columns = df.select_dtypes(include=['object']).columns

        # Step 2: Auto-Generated Dashboard with Key Visualizations
        st.write("## Auto-Generated Dashboard")

        # Distribution plots for numerical columns
        if len(numeric_columns) > 0:
            st.write("### Distribution of Numerical Columns")
            for i in range(0, len(numeric_columns), 2):
                cols = st.columns(2)
                for idx, column in enumerate(numeric_columns[i:i + 2]):
                    with cols[idx]:
                        fig, ax = plt.subplots(figsize=(6, 4))
                        sns.histplot(df[column], kde=True, color='lightblue', bins=30, ax=ax)
                        ax.set_title(f'Distribution of {column}')
                        st.pyplot(fig)

                        # Simple prompt for AI description
                        description_prompt = f"Provide a simple explanation for the distribution of '{column}'. Max 45 words."

                        # Create a button to trigger AI description generation
                        if st.button(f"See Explanation for {column}"):
                            ai_description = generate_ai_description(description_prompt)
                            st.write(f"*AI's Description:* {ai_description}")
                        else:
                            st.write("Click 'See Explanation' to generate an AI description for this plot.")

                        # Download button for the plot
                        st.download_button("Download Plot", data=download_plot(fig, f"{column}_distribution.png"),
                                           file_name=f"{column}_distribution.png")

        # Bar plots for categorical columns
        if len(categorical_columns) > 0:
            st.write("### Distribution of Categorical Columns")
            for i in range(0, len(categorical_columns), 2):
                cols = st.columns(2)
                for idx, column in enumerate(categorical_columns[i:i + 2]):
                    with cols[idx]:
                        fig, ax = plt.subplots(figsize=(6, 4))
                        top_categories = df[column].value_counts().nlargest(10)
                        sns.countplot(x=column, data=df[df[column].isin(top_categories.index)], palette="Set2",
                                      order=top_categories.index, ax=ax)
                        ax.set_title(f'Distribution of {column}')
                        plt.xticks(rotation=45)
                        st.pyplot(fig)

                        # Description prompt for a simple explanation of the bar plot
                        description_prompt = (
                            f"Summarize the bar plot (max 45 words) showing the distribution of '{column}' as if you can see it. "
                            "Describe the general trends one might expect to see in the most frequent categories and any patterns "
                            "that are typically observed in such distributions. Keep it concise and avoid requesting additional data."
                        )

                        # Button to generate AI description
                        if st.button(f"See Explanation for {column}"):
                            ai_description = generate_ai_description(description_prompt)
                            truncated_response = (
                                " ".join(ai_description.split()[:MAX_TOKENS]) + "..."
                                if len(ai_description.split()) > MAX_TOKENS else ai_description
                            )
                            st.write(f"*AI's Description:* {truncated_response}")
                        else:
                            st.write("Click 'See Explanation' to generate an AI description for this plot.")

                        # Download button for the plot
                        st.download_button("Download Plot", data=download_plot(fig, f"{column}_barplot.png"),
                                           file_name=f"{column}_barplot.png")

            # Session State Initialization
            if "fig" not in st.session_state:
                st.session_state["fig"] = None
            if "description" not in st.session_state:
                st.session_state["description"] = None

            # Custom Visualization Section
            st.write("## Generate Custom Visualizations")
            plot_type = st.selectbox("Select Plot Type",
                                     ["Histogram", "Bar Plot", "Line Plot", "Scatter Plot", "Pie Chart",
                                      "Box Plot", "Correlation Heatmap", "Violin Plot", "Pair Plot"])

            # Plot generation logic with AI description
            if plot_type == "Histogram":
                selected_column = st.selectbox("Select column for histogram", numeric_columns)
                if st.button("Generate Histogram"):
                    fig, ax = plt.subplots(figsize=(8, 4))
                    sns.histplot(df[selected_column], bins=30, kde=True, color='lightblue', ax=ax)
                    ax.set_title(f'Histogram of {selected_column}')
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            elif plot_type == "Bar Plot":
                selected_column = st.selectbox("Select column for bar plot", categorical_columns)
                if st.button("Generate Bar Plot"):
                    fig, ax = plt.subplots(figsize=(8, 4))
                    top_categories = df[selected_column].value_counts().nlargest(10)
                    sns.countplot(x=selected_column, data=df[df[selected_column].isin(top_categories.index)],
                                  palette="Set2", order=top_categories.index, ax=ax)
                    ax.set_title(f'Top 10 Categories in {selected_column}')
                    ax.tick_params(axis='x', rotation=45)
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            elif plot_type == "Line Plot":
                x_column = st.selectbox("Select X-axis column for line plot", numeric_columns)
                y_column = st.selectbox("Select Y-axis column for line plot", numeric_columns)
                if st.button("Generate Line Plot"):
                    fig, ax = plt.subplots(figsize=(8, 4))
                    sns.lineplot(data=df, x=x_column, y=y_column, ax=ax)
                    ax.set_title(f'Line Plot of {y_column} vs {x_column}')
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            elif plot_type == "Scatter Plot":
                x_column = st.selectbox("Select X-axis column for scatter plot", numeric_columns)
                y_column = st.selectbox("Select Y-axis column for scatter plot", numeric_columns)
                if st.button("Generate Scatter Plot"):
                    fig, ax = plt.subplots(figsize=(8, 4))
                    sns.scatterplot(data=df, x=x_column, y=y_column, ax=ax)
                    ax.set_title(f'Scatter Plot of {y_column} vs {x_column}')
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            elif plot_type == "Pie Chart":
                selected_column = st.selectbox("Select column for pie chart", categorical_columns)
                if st.button("Generate Pie Chart"):
                    fig, ax = plt.subplots(figsize=(8, 4))
                    df[selected_column].value_counts().plot.pie(autopct='%1.1f%%', colors=sns.color_palette("Set3"),
                                                                ax=ax)
                    ax.set_ylabel('')
                    ax.set_title(f'Pie Chart of {selected_column}')
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            elif plot_type == "Box Plot":
                selected_column = st.selectbox("Select column for box plot", numeric_columns)
                if st.button("Generate Box Plot"):
                    fig, ax = plt.subplots(figsize=(8, 4))
                    sns.boxplot(y=selected_column, data=df, ax=ax, color='lightgreen')
                    ax.set_title(f'Box Plot of {selected_column}')
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            elif plot_type == "Correlation Heatmap":
                if st.button("Generate Correlation Heatmap"):
                    fig, ax = plt.subplots(figsize=(10, 8))
                    sns.heatmap(df[numeric_columns].corr(), annot=True, cmap='coolwarm', ax=ax)
                    ax.set_title('Correlation Heatmap')
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            elif plot_type == "Violin Plot":
                selected_column = st.selectbox("Select column for violin plot", numeric_columns)
                category_column = st.selectbox("Select category column for grouping", categorical_columns)
                if st.button("Generate Violin Plot"):
                    fig, ax = plt.subplots(figsize=(8, 4))
                    sns.violinplot(x=category_column, y=selected_column, data=df, ax=ax, palette="Set2")
                    ax.set_title(f'Violin Plot of {selected_column} by {category_column}')
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            elif plot_type == "Pair Plot":
                if st.button("Generate Pair Plot"):
                    fig = sns.pairplot(df[numeric_columns])
                    st.session_state["fig"] = fig
                    st.session_state["description"] = None

            # Display the plot if it exists in session state
            if st.session_state["fig"] is not None:
                if plot_type == "Pair Plot":
                    st.pyplot(st.session_state["fig"])  # For Pair Plot generated by sns.pairplot
                else:
                    st.pyplot(st.session_state["fig"])

                # AI description generation button and display
                if st.button("Show AI Description"):
                    description_prompt = f"Provide a brief description of the {plot_type.lower()} (max 50 words) for {selected_column if 'selected_column' in locals() else x_column + ' and ' + y_column}."
                    ai_description = generate_ai_description(description_prompt)
                    truncated_response = " ".join(ai_description.split()[:MAX_TOKENS]) + "..." if len(
                        ai_description.split()) > MAX_TOKENS else ai_description
                    st.session_state["description"] = truncated_response

                if st.session_state["description"] is not None:
                    st.write(f"AI's Description: {st.session_state['description']}")

                # Download button for the plot
                file_name = f"{plot_type.replace(' ', '_').lower()}.png"
                st.download_button("Download Plot", data=download_plot(st.session_state["fig"], file_name),
                                   file_name=file_name)

            # AI-based Querying of Dataset
        st.write("## Ask a Question about the Dataset")
        # Assume st.text_input is where the user enters the question
        user_question = st.text_input("Ask a question based on the data")

        # Check if user_question is not empty before creating the prompt
        if user_question:
            data_context = df.head().to_string()  # or summarize the data differently
            prompt = f"Using the following data:\n\n{data_context}\n\nQuestion: {user_question}"

            # Pass this prompt to the AI model
            messages = [HumanMessage(content=prompt)]
            ai_response = chat_llm(messages=messages)
            st.write(f"AI's Answer: {ai_response.content}")
        else:
            st.write("Please enter a question to get an answer.")

        # AI-Based Correlation Insights
        st.write("## AI-Based Correlation Insights")

        if len(numeric_columns) > 1:
            correlation_matrix = df[numeric_columns].corr().to_dict()
            if st.button("Generate AI Insights on Correlations"):
                formatted_corr_matrix = pd.DataFrame(correlation_matrix).to_string()
                ai_response = chat_llm(messages=[HumanMessage(content=formatted_corr_matrix)])
                st.write(f"AI Insights on Correlations: {ai_response.content}")
        else:
            st.write("Not enough numerical columns to compute correlations.")


# Render the appropriate page based on the current session state
if st.session_state.page == 'main':
    show_main_page()
elif st.session_state.page == 'qna_history':
    show_qna_history()
